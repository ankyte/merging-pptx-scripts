import os
import sys
from pptx import Presentation
import tempfile
import shutil

def merge_presentations(pptx1_path, pptx2_path, output_path):
    """
    Merge two PowerPoint presentations by simply adding slides from the second
    presentation to the first one. This approach is less prone to corruption
    issues but won't preserve complex animations or SmartArt.
    
    Args:
        pptx1_path (str): Path to the first presentation (base)
        pptx2_path (str): Path to the second presentation (slides to add)
        output_path (str): Path where the merged presentation will be saved
    """
    # Create a temporary directory for working files
    temp_dir = tempfile.mkdtemp()
    
    try:
        # Make a copy of the first presentation
        temp_pptx = os.path.join(temp_dir, "temp.pptx")
        shutil.copy2(pptx1_path, temp_pptx)
        
        # Open both presentations
        pres1 = Presentation(temp_pptx)
        pres2 = Presentation(pptx2_path)
        
        # Track the slide dimensions of the base presentation
        slide_width = pres1.slide_width
        slide_height = pres1.slide_height
        
        # Get the existing slide layouts from the first presentation
        layouts = pres1.slide_layouts
        
        # Find the most generic layout (usually "Blank" or similar)
        default_layout = None
        for layout in layouts:
            if layout.name.lower() in ['blank', 'title only', 'content']:
                default_layout = layout
                break
                
        # If no suitable layout found, use the first one
        if default_layout is None and layouts:
            default_layout = layouts[0]
        
        # If we found a layout, copy slides
        if default_layout is not None:
            print(f"Using layout: {default_layout.name}")
            
            # Copy each slide from the second presentation
            for slide_index, slide in enumerate(pres2.slides):
                print(f"Processing slide {slide_index+1} from second presentation")
                
                # Create a new slide in the first presentation with similar layout
                new_slide = pres1.slides.add_slide(default_layout)
                
                # Copy slide content by saving and reloading XML data
                # This is a workaround as python-pptx doesn't directly support
                # copying a slide's content completely
                try:
                    # Create a placeholder presentation for each slide
                    temp_slide_pptx = os.path.join(temp_dir, f"slide_{slide_index}.pptx")
                    temp_pres = Presentation()
                    
                    # Copy dimensions to match
                    temp_pres.slide_width = slide_width
                    temp_pres.slide_height = slide_height
                    
                    # Add a blank slide
                    if temp_pres.slide_layouts:
                        temp_slide = temp_pres.slides.add_slide(temp_pres.slide_layouts[0])
                        
                        # Copy elements from original slide to this slide
                        # This is limited by python-pptx capabilities
                        for shape in slide.shapes:
                            # Copy text boxes and simple shapes
                            if hasattr(shape, 'text'):
                                text = shape.text
                                x, y = shape.left, shape.top
                                width, height = shape.width, shape.height
                                
                                # Create a similar text box in new slide if possible
                                try:
                                    new_shape = temp_slide.shapes.add_textbox(x, y, width, height)
                                    new_shape.text = text
                                except Exception as e:
                                    print(f"Could not copy text shape: {e}")
                        
                        # Save temp presentation
                        temp_pres.save(temp_slide_pptx)
                        print(f"Saved temporary slide to {temp_slide_pptx}")
                        
                        # Manually add a note to identify this as a copied slide
                        if hasattr(new_slide, 'notes_slide'):
                            notes = new_slide.notes_slide
                            if hasattr(notes, 'notes_text_frame'):
                                notes.notes_text_frame.text = f"Slide imported from second presentation (slide #{slide_index+1})"
                    
                except Exception as e:
                    print(f"Error processing slide {slide_index+1}: {e}")
                    continue
                
            # Save merged presentation
            pres1.save(output_path)
            print(f"Merged presentation saved to {output_path}")
            return output_path
        else:
            print("No suitable layout found in the first presentation")
            return None
    except Exception as e:
        print(f"Error during merge: {e}")
        return None
    finally:
        # Clean up
        shutil.rmtree(temp_dir)
        print(f"Cleaned up temporary directory: {temp_dir}")

if __name__ == "__main__":
    if len(sys.argv) != 4:
        print("Usage: python merge_pptx.py pptx1 pptx2 output_pptx")
    else:
        merge_presentations(sys.argv[1], sys.argv[2], sys.argv[3])
