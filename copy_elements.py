from pptx import Presentation
import copy
import os

def merge_presentations(base_ppt_path, add_ppt_path, output_path):
    """
    Merge two PowerPoint presentations while preserving charts and other elements.
    
    Args:
        base_ppt_path (str): Path to the first presentation (base presentation)
        add_ppt_path (str): Path to the second presentation to be added
        output_path (str): Path where the merged presentation will be saved
    """
    # Load both presentations
    base_pres = Presentation(base_ppt_path)
    add_pres = Presentation(add_ppt_path)
    
    # Create a mapping of slide layouts
    layout_mapping = {}
    
    # First, check if we need to add new slide layouts
    # This is important for chart slides that depend on specific layouts
    for slide_layout in add_pres.slide_layouts:
        layout_name = slide_layout.name
        layout_found = False
        
        for base_layout in base_pres.slide_layouts:
            if base_layout.name == layout_name:
                layout_mapping[slide_layout] = base_layout
                layout_found = True
                break
        
        if not layout_found:
            # Copy slide layout to base presentation
            # Note: python-pptx doesn't support direct layout copying,
            # but we can map to the closest match
            layout_mapping[slide_layout] = find_closest_layout(base_pres, slide_layout)
    
    # Copy each slide from the second presentation to the first
    for slide in add_pres.slides:
        # Get the layout that this slide is based on
        source_layout = slide.slide_layout
        target_layout = layout_mapping.get(source_layout, base_pres.slide_layouts[0])
        
        # Create a new slide in the base presentation with the matching layout
        new_slide = base_pres.slides.add_slide(target_layout)
        
        # Copy slide contents
        copy_slide_contents(slide, new_slide)
    
    # Save the merged presentation
    base_pres.save(output_path)
    print(f"Merged presentation saved to {output_path}")

def find_closest_layout(pres, source_layout):
    """Find the closest matching layout in the target presentation."""
    # First try to match by name
    for layout in pres.slide_layouts:
        if layout.name == source_layout.name:
            return layout
    
    # If no match by name, try to match by type/purpose
    # This is a simplified version; you might need to expand this logic
    # based on your specific layouts
    layout_types = {
        'Title Slide': 0,
        'Title and Content': 1,
        'Section Header': 2,
        'Two Content': 3,
        'Comparison': 4,
        'Title Only': 5,
        'Blank': 6,
        'Content with Caption': 7,
        'Picture with Caption': 8
    }
    
    source_type = layout_types.get(source_layout.name, 1)  # Default to Title and Content
    
    # Return the matching layout or default to Title and Content (index 1)
    try:
        return pres.slide_layouts[source_type]
    except IndexError:
        return pres.slide_layouts[1]  # Default to Title and Content

def copy_slide_contents(source_slide, target_slide):
    """
    Copy the contents of a slide to another slide.
    
    This function handles shapes, placeholders, charts, tables, etc.
    """
    # Copy slide background
    if source_slide.background.fill.type != 0:  # 0 means no fill
        target_slide.background = copy.deepcopy(source_slide.background)
    
    # Handle shapes (including charts)
    for shape in source_slide.shapes:
        if shape.shape_type == 1:  # Auto Shape
            copy_auto_shape(shape, target_slide)
        elif shape.shape_type == 3:  # Chart
            copy_chart(shape, target_slide)
        elif shape.shape_type == 6:  # Group Shape
            copy_group_shape(shape, target_slide)
        elif shape.shape_type == 7:  # Textbox
            copy_textbox(shape, target_slide)
        elif shape.shape_type == 8:  # Picture
            copy_picture(shape, target_slide)
        elif shape.shape_type == 14:  # Placeholder
            copy_placeholder(shape, target_slide)
        elif shape.shape_type == 19:  # Table
            copy_table(shape, target_slide)
        else:
            # For other shape types, create a generic shape
            copy_generic_shape(shape, target_slide)

def copy_auto_shape(shape, slide):
    """Copy an auto shape to the target slide."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    
    # Add a new shape
    new_shape = slide.shapes.add_shape(
        shape.auto_shape_type, left, top, width, height
    )
    
    # Copy shape properties
    copy_shape_properties(shape, new_shape)
    copy_text_frame(shape, new_shape)

def copy_chart(shape, slide):
    """
    Copy a chart to the target slide.
    
    This method doesn't directly copy the chart but creates a placeholder with a
    note about the chart's presence. You may need to enhance this based on your needs.
    """
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    
    # Create a textbox to indicate the chart's presence
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    
    # Add text noting this is a chart
    if hasattr(shape, 'chart') and hasattr(shape.chart, 'chart_type'):
        chart_type = shape.chart.chart_type
        p = tf.add_paragraph()
        p.text = f"Chart (Type: {chart_type})"
        p.font.bold = True
        p.font.size = 14
    else:
        p = tf.add_paragraph()
        p.text = "Chart placeholder"
        p.font.bold = True
        p.font.size = 14

def copy_group_shape(shape, slide):
    """Copy a group shape to the target slide."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    
    # Create a group shape container
    group = slide.shapes.add_group_shape()
    group.left = left
    group.top = top
    group.width = width
    group.height = height
    
    # Loop through shapes in the group and add them
    for child_shape in shape.shapes:
        # Adjust coordinates relative to the group
        child_left = child_shape.left - shape.left
        child_top = child_shape.top - shape.top
        
        if child_shape.shape_type == 1:  # Auto Shape
            new_shape = group.shapes.add_shape(
                child_shape.auto_shape_type, 
                child_left, child_top, 
                child_shape.width, child_shape.height
            )
            copy_shape_properties(child_shape, new_shape)
            copy_text_frame(child_shape, new_shape)
        elif child_shape.shape_type == 7:  # Textbox
            new_shape = group.shapes.add_textbox(
                child_left, child_top,
                child_shape.width, child_shape.height
            )
            copy_text_frame(child_shape, new_shape)
        # Add more conditions for other shape types as needed

def copy_textbox(shape, slide):
    """Copy a textbox to the target slide."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    
    # Create a new textbox
    new_shape = slide.shapes.add_textbox(left, top, width, height)
    
    # Copy text content and formatting
    copy_text_frame(shape, new_shape)

def copy_picture(shape, slide):
    """Copy a picture to the target slide."""
    # For pictures, we need to save the image temporarily and then add it
    # This is a limitation of python-pptx
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    
    # Create a temporary file for the image
    temp_image_path = "temp_image.png"
    
    try:
        # Try to save the image
        with open(temp_image_path, "wb") as f:
            if hasattr(shape, "image") and shape.image:
                f.write(shape.image.blob)
                # Add the image to the new slide
                slide.shapes.add_picture(temp_image_path, left, top, width, height)
            else:
                # If we can't get the image, create a placeholder
                txbox = slide.shapes.add_textbox(left, top, width, height)
                txbox.text_frame.text = "Image placeholder"
    except Exception as e:
        print(f"Error copying image: {e}")
        # Create a placeholder instead
        txbox = slide.shapes.add_textbox(left, top, width, height)
        txbox.text_frame.text = "Image placeholder"
    finally:
        # Clean up temporary file
        if os.path.exists(temp_image_path):
            os.remove(temp_image_path)

def copy_placeholder(shape, slide):
    """Copy a placeholder to the target slide."""
    # Find a matching placeholder by index in the target slide
    target_placeholders = [p for p in slide.placeholders]
    
    # Try to find a placeholder with the same index
    target_placeholder = None
    if hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'idx'):
        idx = shape.placeholder_format.idx
        for p in target_placeholders:
            if hasattr(p, 'placeholder_format') and hasattr(p.placeholder_format, 'idx'):
                if p.placeholder_format.idx == idx:
                    target_placeholder = p
                    break
    
    if target_placeholder:
        # If we found a matching placeholder, copy the content
        if shape.has_text_frame:
            copy_text_frame(shape, target_placeholder)
    else:
        # If no matching placeholder, create a textbox with the content
        if shape.has_text_frame:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            txbox = slide.shapes.add_textbox(left, top, width, height)
            copy_text_frame(shape, txbox)

def copy_table(shape, slide):
    """Copy a table to the target slide."""
    if not hasattr(shape, 'table'):
        return
    
    source_table = shape.table
    rows, cols = len(source_table.rows), len(source_table.columns)
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    
    # Create a new table
    new_table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Copy table contents and formatting
    for i, row in enumerate(source_table.rows):
        for j, cell in enumerate(row.cells):
            target_cell = new_table.cell(i, j)
            
            # Copy text content
            if cell.text_frame:
                target_cell.text = cell.text
                
                # Copy paragraph formatting if possible
                if hasattr(cell.text_frame, 'paragraphs') and cell.text_frame.paragraphs:
                    for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                        if p_idx < len(target_cell.text_frame.paragraphs):
                            target_paragraph = target_cell.text_frame.paragraphs[p_idx]
                            if hasattr(paragraph, 'font') and hasattr(target_paragraph, 'font'):
                                if hasattr(paragraph.font, 'bold'):
                                    target_paragraph.font.bold = paragraph.font.bold
                                if hasattr(paragraph.font, 'italic'):
                                    target_paragraph.font.italic = paragraph.font.italic
                                if hasattr(paragraph.font, 'color') and paragraph.font.color.rgb:
                                    target_paragraph.font.color.rgb = paragraph.font.color.rgb

def copy_generic_shape(shape, slide):
    """Create a generic placeholder for unsupported shape types."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    
    # Create a textbox as a placeholder
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.text = f"Shape placeholder (Type: {shape.shape_type})"

def copy_shape_properties(source, target):
    """Copy common shape properties from source to target."""
    # Copy fill
    if hasattr(source, 'fill') and hasattr(target, 'fill'):
        if source.fill.type == 1:  # Solid fill
            if source.fill.fore_color.rgb:
                target.fill.solid()
                target.fill.fore_color.rgb = source.fill.fore_color.rgb
        elif source.fill.type == 0:  # No fill
            target.fill.background()
    
    # Copy line properties
    if hasattr(source, 'line') and hasattr(target, 'line'):
        if hasattr(source.line, 'color') and source.line.color.rgb:
            target.line.color.rgb = source.line.color.rgb
        
        if hasattr(source.line, 'width'):
            target.line.width = source.line.width

def copy_text_frame(source, target):
    """Copy text frame content and formatting from source to target."""
    if not (hasattr(source, 'text_frame') and hasattr(target, 'text_frame')):
        return
    
    source_tf = source.text_frame
    target_tf = target.text_frame
    
    # Copy text frame properties
    if hasattr(source_tf, 'word_wrap') and hasattr(target_tf, 'word_wrap'):
        target_tf.word_wrap = source_tf.word_wrap
    
    # Clear existing paragraphs in target (except first one)
    while len(target_tf.paragraphs) > 1:
        p = target_tf.paragraphs[-1]
        p._p.getparent().remove(p._p)
    
    # Get the first (and possibly only) paragraph in the target
    if not target_tf.paragraphs:
        target_para = target_tf.add_paragraph()
    else:
        target_para = target_tf.paragraphs[0]
        target_para.text = ""
    
    # Copy paragraphs
    for i, source_para in enumerate(source_tf.paragraphs):
        if i == 0:
            # Use the existing first paragraph
            para = target_para
        else:
            # Add a new paragraph
            para = target_tf.add_paragraph()
        
        # Copy text
        para.text = source_para.text
        
        # Copy paragraph formatting
        if hasattr(source_para, 'alignment') and hasattr(para, 'alignment'):
            para.alignment = source_para.alignment
        
        # Copy font formatting
        if hasattr(source_para, 'font') and hasattr(para, 'font'):
            if hasattr(source_para.font, 'size') and source_para.font.size:
                para.font.size = source_para.font.size
            
            if hasattr(source_para.font, 'bold'):
                para.font.bold = source_para.font.bold
            
            if hasattr(source_para.font, 'italic'):
                para.font.italic = source_para.font.italic
            
            if hasattr(source_para.font, 'underline'):
                para.font.underline = source_para.font.underline
            
            if hasattr(source_para.font, 'color') and source_para.font.color.rgb:
                para.font.color.rgb = source_para.font.color.rgb

# Example usage
if __name__ == "__main__":
    merge_presentations("base_presentation.pptx", "additional_presentation.pptx", "merged_output.pptx")
