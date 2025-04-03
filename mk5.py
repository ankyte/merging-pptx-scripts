#!/usr/bin/env python3
"""
PPT Merger - Merge two PowerPoint presentations with the same slide layouts.
This script runs on Unix and handles various PowerPoint elements including charts, text, and shapes.

Usage: python ppt_merger.py <input_ppt1> <input_ppt2> <output_ppt>
"""

import os
import sys
import copy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import argparse
import logging

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def validate_files(input_files, output_file):
    """Validate input and output files."""
    # Check if input files exist
    for file in input_files:
        if not os.path.exists(file):
            raise FileNotFoundError(f"Input file not found: {file}")
        
        # Check if input files are PowerPoint files
        if not file.lower().endswith(('.pptx', '.ppt')):
            raise ValueError(f"Input file is not a PowerPoint file: {file}")
    
    # Check if output directory exists
    output_dir = os.path.dirname(output_file)
    if output_dir and not os.path.exists(output_dir):
        raise FileNotFoundError(f"Output directory does not exist: {output_dir}")
    
    # Check if output file is a PowerPoint file
    if not output_file.lower().endswith(('.pptx', '.ppt')):
        raise ValueError(f"Output file is not a PowerPoint file: {output_file}")

def copy_slide_layouts(source_pres, target_pres):
    """Copy slide layouts from source presentation to target presentation if needed."""
    source_layouts = {layout.name: layout for layout in source_pres.slide_layouts}
    target_layouts = {layout.name: layout for layout in target_pres.slide_layouts}
    
    # Check if any layouts are missing in the target
    missing_layouts = set(source_layouts.keys()) - set(target_layouts.keys())
    if missing_layouts:
        logger.warning(f"Some layouts might be missing in the merged presentation: {missing_layouts}")
    
    return source_layouts, target_layouts

def copy_slide(source_slide, target_pres, source_layouts, target_layouts):
    """Copy a slide from source presentation to target presentation."""
    # Find the closest matching layout
    if source_slide.slide_layout.name in target_layouts:
        target_layout = target_layouts[source_slide.slide_layout.name]
    else:
        # Fallback to a default layout if the exact match is not found
        logger.warning(f"Layout '{source_slide.slide_layout.name}' not found in target. Using default.")
        target_layout = target_pres.slide_layouts[0]
    
    # Create a new slide with the matched layout
    target_slide = target_pres.slides.add_slide(target_layout)
    
    # Copy slide properties
    if hasattr(source_slide, 'background'):
        if hasattr(target_slide, 'background'):
            try:
                target_slide.background = copy.deepcopy(source_slide.background)
            except Exception as e:
                logger.warning(f"Could not copy slide background: {e}")
    
    # Copy shapes from source to target
    for shape in source_slide.shapes:
        try:
            copy_shape(shape, target_slide)
        except Exception as e:
            logger.warning(f"Error copying shape: {e}")
    
    return target_slide

def copy_shape(shape, target_slide):
    """Copy a shape from source slide to target slide."""
    # Handle different shape types
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        copy_placeholder(shape, target_slide)
    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
        copy_chart(shape, target_slide)
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        copy_table(shape, target_slide)
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        copy_picture(shape, target_slide)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        copy_group(shape, target_slide)
    else:
        copy_generic_shape(shape, target_slide)

def copy_placeholder(shape, target_slide):
    """Copy a placeholder shape."""
    # Find a matching placeholder in the target slide
    for placeholder in target_slide.placeholders:
        if placeholder.placeholder_format.type == shape.placeholder_format.type:
            # Copy text from source placeholder to target placeholder
            if hasattr(shape, 'text') and shape.text:
                placeholder.text = shape.text
            
            # Copy formatting if available
            if hasattr(shape, 'text_frame') and hasattr(placeholder, 'text_frame'):
                copy_text_frame(shape.text_frame, placeholder.text_frame)
            
            break

def copy_text_frame(source_frame, target_frame):
    """Copy text frame including formatting."""
    if not hasattr(source_frame, 'paragraphs'):
        return
    
    # Clear existing paragraphs in target
    while len(target_frame.paragraphs) > 1:
        p = target_frame.paragraphs[-1]
        p._p.getparent().remove(p._p)
    
    # Copy each paragraph
    for i, source_para in enumerate(source_frame.paragraphs):
        if i == 0 and len(target_frame.paragraphs) > 0:
            target_para = target_frame.paragraphs[0]
        else:
            target_para = target_frame.add_paragraph()
        
        # Copy text
        target_para.text = source_para.text
        
        # Copy paragraph formatting
        if hasattr(source_para, 'alignment') and source_para.alignment:
            target_para.alignment = source_para.alignment
        
        if hasattr(source_para, 'level') and source_para.level:
            target_para.level = source_para.level
        
        # Copy runs (text formatting)
        if hasattr(source_para, 'runs'):
            for j, source_run in enumerate(source_para.runs):
                if j == 0 and len(target_para.runs) > 0:
                    target_run = target_para.runs[0]
                else:
                    target_run = target_para.add_run()
                
                target_run.text = source_run.text
                
                # Copy run formatting
                if hasattr(source_run, 'font') and hasattr(target_run, 'font'):
                    if hasattr(source_run.font, 'bold'):
                        target_run.font.bold = source_run.font.bold
                    if hasattr(source_run.font, 'italic'):
                        target_run.font.italic = source_run.font.italic
                    if hasattr(source_run.font, 'underline'):
                        target_run.font.underline = source_run.font.underline
                    if hasattr(source_run.font, 'size'):
                        target_run.font.size = source_run.font.size
                    if hasattr(source_run.font, 'color'):
                        target_run.font.color.rgb = source_run.font.color.rgb

def copy_chart(shape, target_slide):
    """Copy a chart shape."""
    # Charts are complex objects in PowerPoint
    # This is a simplified approach that creates a placeholder noting there was a chart
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    chart_placeholder = target_slide.shapes.add_textbox(left, top, width, height)
    chart_placeholder.text = f"[Chart: {shape.name if hasattr(shape, 'name') else 'Unknown'}]"
    chart_placeholder.text_frame.paragraphs[0].font.bold = True
    chart_placeholder.text_frame.paragraphs[0].alignment = 1  # Center

def copy_table(shape, target_slide):
    """Copy a table shape."""
    if not hasattr(shape, 'table'):
        return
    
    # Get table dimensions
    rows, cols = len(shape.table.rows), len(shape.table.columns)
    
    # Create a new table
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    table = target_slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Copy cell contents
    for i, row in enumerate(shape.table.rows):
        for j, cell in enumerate(row.cells):
            if i < len(table.rows) and j < len(table.rows[i].cells):
                table.rows[i].cells[j].text = cell.text

def copy_picture(shape, target_slide):
    """Copy a picture shape."""
    # Pictures require special handling which is beyond the scope of this script
    # Create a placeholder instead
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    pic_placeholder = target_slide.shapes.add_textbox(left, top, width, height)
    pic_placeholder.text = f"[Image: {shape.name if hasattr(shape, 'name') else 'Unknown'}]"
    pic_placeholder.text_frame.paragraphs[0].font.bold = True
    pic_placeholder.text_frame.paragraphs[0].alignment = 1  # Center

def copy_group(shape, target_slide):
    """Copy a group shape."""
    # Group shapes contain multiple shapes
    # Create a placeholder for the group
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    group_placeholder = target_slide.shapes.add_textbox(left, top, width, height)
    group_placeholder.text = f"[Group: {shape.name if hasattr(shape, 'name') else 'Group of shapes'}]"
    group_placeholder.text_frame.paragraphs[0].font.bold = True
    group_placeholder.text_frame.paragraphs[0].alignment = 1  # Center

def copy_generic_shape(shape, target_slide):
    """Copy a generic shape."""
    try:
        # Try to add a similar shape (this is simplified)
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        
        # If the shape has text, create a textbox with the same text
        if hasattr(shape, 'text') and shape.text:
            textbox = target_slide.shapes.add_textbox(left, top, width, height)
            textbox.text = shape.text
            
            # Copy text formatting if possible
            if hasattr(shape, 'text_frame') and hasattr(textbox, 'text_frame'):
                copy_text_frame(shape.text_frame, textbox.text_frame)
        else:
            # For shapes without text, add a rectangle as a placeholder
            rect = target_slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
            if hasattr(shape, 'fill') and hasattr(rect, 'fill'):
                if hasattr(shape.fill, 'fore_color') and hasattr(rect.fill, 'fore_color'):
                    rect.fill.fore_color.rgb = shape.fill.fore_color.rgb
    
    except Exception as e:
        logger.warning(f"Could not copy generic shape: {e}")

def merge_presentations(input_files, output_file):
    """Merge multiple PowerPoint presentations into one."""
    # Validate files
    validate_files(input_files, output_file)
    
    # Create a new presentation for output
    merged_pres = Presentation(input_files[0])
    
    # Process each input presentation (skip the first one as it's already the base)
    for i, input_file in enumerate(input_files[1:], 1):
        logger.info(f"Processing presentation {i+1}: {input_file}")
        
        # Open the source presentation
        source_pres = Presentation(input_file)
        
        # Copy slide layouts if needed
        source_layouts, target_layouts = copy_slide_layouts(source_pres, merged_pres)
        
        # Copy each slide
        for j, slide in enumerate(source_pres.slides):
            logger.info(f"  Copying slide {j+1}")
            copy_slide(slide, merged_pres, source_layouts, target_layouts)
    
    # Save the merged presentation
    logger.info(f"Saving merged presentation to: {output_file}")
    merged_pres.save(output_file)
    logger.info("Merge completed successfully!")

def main():
    """Main function to process command line arguments and merge presentations."""
    parser = argparse.ArgumentParser(description='Merge PowerPoint presentations with the same slide layouts.')
    parser.add_argument('input_files', nargs='+', help='Input PowerPoint files')
    parser.add_argument('output_file', help='Output PowerPoint file')
    parser.add_argument('--debug', action='store_true', help='Enable debug logging')
    
    args = parser.parse_args()
    
    if args.debug:
        logger.setLevel(logging.DEBUG)
    
    try:
        merge_presentations(args.input_files[:-1], args.input_files[-1])
    except Exception as e:
        logger.error(f"Error merging presentations: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
